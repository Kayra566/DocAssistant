from __future__ import annotations

from io import BytesIO

from app.core.config import settings
from app.core.exceptions import ValidationError

# Her eleman (sayfa_no, metin) — sayfa 1'den başlar.
Pages = list[tuple[int, str]]


def _extract_pdf(data: bytes) -> Pages:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    pages: Pages = []
    for i, page in enumerate(reader.pages, start=1):
        pages.append((i, page.extract_text() or ""))
    return pages


def _extract_docx(data: bytes) -> Pages:
    from docx import Document as DocxDocument

    doc = DocxDocument(BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs if p.text)
    return [(1, text)]


def _extract_xlsx(data: bytes) -> Pages:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    pages: Pages = []
    for idx, ws in enumerate(wb.worksheets, start=1):
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        pages.append((idx, "\n".join(rows)))
    wb.close()
    return pages or [(1, "")]


def _extract_pptx(data: bytes) -> Pages:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    pages: Pages = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text
        ]
        pages.append((idx, "\n".join(texts)))
    return pages or [(1, "")]


def _extract_text(data: bytes) -> Pages:
    return [(1, data.decode("utf-8", errors="replace"))]


def _extract_image(data: bytes) -> Pages:
    if not settings.ENABLE_OCR:
        raise ValidationError("OCR devre dışı (ENABLE_OCR=false).")
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise ValidationError("OCR bağımlılıkları kurulu değil.") from exc
    text = pytesseract.image_to_string(Image.open(BytesIO(data)))
    return [(1, text)]


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "xlsx": _extract_xlsx,
    "pptx": _extract_pptx,
    "txt": _extract_text,
    "md": _extract_text,
    "image": _extract_image,
}


def extract_pages(file_type: str, data: bytes) -> Pages:
    extractor = _EXTRACTORS.get(file_type)
    if not extractor:
        raise ValidationError(f"Bu tip için metin çıkarma yok: {file_type}")
    return extractor(data)
